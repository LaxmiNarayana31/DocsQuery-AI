import tempfile
import traceback

import docx
import pptx
import pandas as pd
from PyPDF2 import PdfReader
from odf.opendocument import load
from odf.text import P
from striprtf.striprtf import rtf_to_text
from spire.presentation import Presentation, FileFormat
from spire.doc import Document as SpireDocument
from spire.doc import FileFormat as SpireFileFormat
from spire.xls import Workbook as SpireWorkbook
from spire.xls import ExcelVersion


class DocumentHelper:
    """Helper class to extract text from multiple document formats."""

    # ---------------- PDF ----------------
    @staticmethod
    def extract_pdf(file):
        pdf_reader = PdfReader(file)
        return "\n".join(
            page.extract_text() or "" for page in pdf_reader.pages
        )

    # ---------------- DOC ----------------
    @staticmethod
    def extract_doc(file):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp_file:
            tmp_file.write(file.read())
            tmp_file_path = tmp_file.name

        document = SpireDocument()
        document.LoadFromFile(tmp_file_path)
        document.SaveToFile(tmp_file_path, SpireFileFormat.Docx2019)
        document.Close()

        doc = docx.Document(tmp_file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    # ---------------- DOCX ----------------
    @staticmethod
    def extract_docx(file):
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])

    # ---------------- TXT ----------------
    @staticmethod
    def extract_txt(file):
        return file.read().decode("utf-8")

    # ---------------- PPT ----------------
    @staticmethod
    def extract_ppt(file):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp_file:
            tmp_file.write(file.read())
            tmp_file_path = tmp_file.name

        pre = Presentation()
        pre.LoadFromFile(tmp_file_path)
        pre.SaveToFile(tmp_file_path, FileFormat.Pptx2019)
        pre.Dispose()

        return DocumentHelper.extract_pptx(open(tmp_file_path, "rb"))

    # ---------------- PPTX ----------------
    @staticmethod
    def extract_pptx(file):
        presentation = pptx.Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    # ---------------- ODT ----------------
    @staticmethod
    def extract_odt(file):
        odt_doc = load(file)
        all_paragraphs = odt_doc.getElementsByType(P)
        text = ""
        for paragraph in all_paragraphs:
            for node in paragraph.childNodes:
                if node.nodeType == node.TEXT_NODE:
                    text += node.data
                elif node.nodeType == node.ELEMENT_NODE and node.tagName == "text:span":
                    for child_node in node.childNodes:
                        if child_node.nodeType == child_node.TEXT_NODE:
                            text += child_node.data
            text += "\n"
        return text

    # ---------------- RTF ----------------
    @staticmethod
    def extract_rtf(file):
        rtf_content = file.read()
        return rtf_to_text(rtf_content.decode("latin-1"))

    # ---------------- CSV ----------------
    @staticmethod
    def extract_csv(file):
        df = pd.read_csv(file)
        rows = [
            ", ".join(f"{col}: {row[col]}" for col in df.columns)
            for _, row in df.iterrows()
        ]
        return "\n".join(rows)

    # ---------------- XLS ----------------
    @staticmethod
    def extract_xls(file):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xls") as tmp_file:
            tmp_file.write(file.read())
            tmp_file_path = tmp_file.name

        workbook = SpireWorkbook()
        workbook.LoadFromFile(tmp_file_path)
        workbook.SaveToFile(tmp_file_path, ExcelVersion.Version2016)
        workbook.Dispose()

        df = pd.read_excel(tmp_file_path)
        return "\n".join(
            [", ".join([f"{col}: {str(row[col])}" for col in df.columns]) for _, row in df.iterrows()]
        )

    # ---------------- XLSX ----------------
    @staticmethod
    def extract_xlsx(file):
        df = pd.read_excel(file)
        return "\n".join(
            [", ".join([f"{col}: {str(row[col])}" for col in df.columns]) for _, row in df.iterrows()]
        )

    # ---------------- Dispatcher ----------------
    @staticmethod
    def extractText(file):
        """
        Extract text from different document formats.
        Supported: pdf, doc, docx, txt, ppt, pptx, odt, rtf, csv, xls, xlsx
        """
        try:
            file_name = file.name.lower()
            ext = file_name.split(".")[-1]

            extractor_map = {
                "pdf": DocumentHelper.extract_pdf,
                "doc": DocumentHelper.extract_doc,
                "docx": DocumentHelper.extract_docx,
                "txt": DocumentHelper.extract_txt,
                "ppt": DocumentHelper.extract_ppt,
                "pptx": DocumentHelper.extract_pptx,
                "odt": DocumentHelper.extract_odt,
                "rtf": DocumentHelper.extract_rtf,
                "csv": DocumentHelper.extract_csv,
                "xls": DocumentHelper.extract_xls,
                "xlsx": DocumentHelper.extract_xlsx,
            }

            if ext not in extractor_map:
                return "Unsupported file format."

            return extractor_map[ext](file)

        except Exception as e:
            traceback_str = traceback.format_exc()
            print(traceback_str)
            line_no = traceback.extract_tb(e.__traceback__)[-1][1]
            print(f"Exception occurred on line {line_no}")
            return str(e)
